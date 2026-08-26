import os
from pptx import Presentation
from pptx.util import Inches

def create_centered_presentation(image_folder, output_file):
    prs = Presentation()
    
    # Optional: Set to Standard 4:3 (10 x 7.5 inches)
    # Default is 16:9. Comment these out if you prefer Widescreen.
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    valid_extensions = ('.png', '.jpg', '.jpeg', '.gif', '.bmp')
    images = [f for f in os.listdir(image_folder) if f.lower().endswith(valid_extensions)]
    images.sort()

    blank_slide_layout = prs.slide_layouts[6]

    for image_name in images:
        slide = prs.slides.add_slide(blank_slide_layout)
        img_path = os.path.join(image_folder, image_name)

        # 1. Add the image to the slide first to get its dimensions
        # We start by placeholder-adding it at 0,0
        pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0))

        # 2. Rescale image to fit the slide while maintaining aspect ratio
        # We calculate the scaling factor for both width and height
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        ratio = min(slide_width / pic.width, slide_height / pic.height)
        
        # pic.width = int(pic.width * ratio)
        # pic.height = int(pic.height * ratio)
        pic.width = Inches(9.5)
        pic.height = Inches(4)

        # 3. Calculate centering positions
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = int((slide_height - pic.height) / 2)
        
        print(f"Centered: {image_name}")

    prs.save(output_file)
    print(f"\nDone! Saved as {output_file}")

if __name__ == "__main__":
    create_centered_presentation('TF-conformation/deeppbs_pdbs/monomer_chains/rmsd/plots', 'centered_presentation.pptx')
