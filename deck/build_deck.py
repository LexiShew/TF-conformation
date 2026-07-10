
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

PAL = {
 'blue':'8ecae6','blue_dk':'4a8fb3','lav':'b5a7e6','lav_dk':'7b68c4',
 'mint':'a8e6cf','mint_dk':'5bbf97','pink':'f6c6d9','pink_dk':'d17aa0',
 'peach':'ffd8be','peach_dk':'e08b5a','ink':'33333a','grid':'e6e6ee',
}
def C(h): return RGBColor.from_string(h)
INK=C(PAL['ink']); ACCENT=C(PAL['lav_dk']); BAND=C('f4f1fb')

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def bg(slide, hexcol):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(hexcol)

def title_bar(slide, title, accent=ACCENT):
    # accent rule + title
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12.1), Inches(0.95))
    tf = box.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=title
    r.font.size=Pt(28); r.font.bold=True; r.font.color.rgb=INK; r.font.name='Calibri'
    # accent underline
    ln = slide.shapes.add_shape(1, Inches(0.62), Inches(1.18), Inches(2.4), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb=accent; ln.line.fill.background()
    return box

def img_fit(slide, path, x, y, maxw, maxh):
    iw,ih = Image.open(path).size
    ar = iw/ih; boxar = maxw/maxh
    if ar>boxar: w=maxw; h=maxw/ar
    else: h=maxh; w=maxh*ar
    x2 = x + (maxw-w)/2; y2 = y + (maxh-h)/2
    return slide.shapes.add_picture(path, int(x2), int(y2), int(w), int(h))

def bullets_body(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.9), Inches(5.6))
    tf=box.text_frame; tf.word_wrap=True
    for i,b in enumerate(bullets):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        run=p.add_run(); run.text=b
        run.font.size=Pt(18); run.font.color.rgb=INK; run.font.name='Calibri'
        p.space_after=Pt(14); p.line_spacing=1.08
        # pastel bullet dot via leading char
        run0=p.runs[0]; run0.text="•  "+run0.text
    return box

prs = Presentation()
prs.slide_width=SW; prs.slide_height=SH
blank = prs.slide_layouts[6]
spec = json.load(open('handoff/deck_spec.json'))

for s in spec:
    slide = prs.slides.add_slide(blank)
    lay=s['layout']
    if lay=='title':
        bg(slide,'f4f1fb')
        # big centered title
        tb=slide.shapes.add_textbox(Inches(0.8),Inches(2.4),Inches(11.7),Inches(1.8))
        tf=tb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=s['title']; r.font.size=Pt(40); r.font.bold=True; r.font.color.rgb=INK
        # accent bar
        ln=slide.shapes.add_shape(1, Inches(5.17),Inches(4.35),Inches(3.0),Pt(4))
        ln.fill.solid(); ln.fill.fore_color.rgb=ACCENT; ln.line.fill.background()
        sb=slide.shapes.add_textbox(Inches(1.2),Inches(4.6),Inches(10.9),Inches(1.8))
        stf=sb.text_frame; stf.word_wrap=True
        for i,line in enumerate(s['subtitle'].split('\n')):
            pp=stf.paragraphs[0] if i==0 else stf.add_paragraph()
            pp.alignment=PP_ALIGN.CENTER; rr=pp.add_run(); rr.text=line
            rr.font.size=Pt(17); rr.font.color.rgb=C(PAL['lav_dk'])
    elif lay=='bullets':
        bg(slide,'ffffff'); title_bar(slide,s['title']); bullets_body(slide,s['bullets'])
    elif lay=='image':
        bg(slide,'ffffff'); title_bar(slide,s['title'])
        img_fit(slide,s['image'],Inches(0.6),Inches(1.4),Inches(12.1),Inches(5.85))
    elif lay=='image_caption':
        bg(slide,'ffffff'); title_bar(slide,s['title'])
        img_fit(slide,s['image'],Inches(0.6),Inches(1.4),Inches(12.1),Inches(5.25))
        cb=slide.shapes.add_textbox(Inches(0.8),Inches(6.75),Inches(11.7),Inches(0.6))
        ctf=cb.text_frame; ctf.word_wrap=True; cp=ctf.paragraphs[0]; cp.alignment=PP_ALIGN.CENTER
        cr=cp.add_run(); cr.text=s['caption']; cr.font.size=Pt(12); cr.font.italic=True; cr.font.color.rgb=C('6a6a75')
    elif lay=='image_grid':
        bg(slide,'ffffff'); title_bar(slide,s['title'])
        imgs=s['images']; 
        # 2x2 grid
        cw,ch=Inches(6.0),Inches(2.85); x0,y0=Inches(0.55),Inches(1.45); gx,gy=Inches(0.25),Inches(0.15)
        for k,ip in enumerate(imgs[:4]):
            rr,cc=divmod(k,2)
            img_fit(slide,ip,int(x0+cc*(cw+gx)),int(y0+rr*(ch+gy)),int(cw),int(ch))
    add_notes(slide, s['notes'])

prs.save('TF_conformation_deck.pptx')
print("saved TF_conformation_deck.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
